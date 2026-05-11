from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- RENK PALETI (Eymen sunumlari icin) ---
MAVI_KOYU = RGBColor(0x1B, 0x3A, 0x5C)
MAVI_ACIK = RGBColor(0x2E, 0x86, 0xC1)
MAVI_CANLI = RGBColor(0x00, 0x96, 0xD6)
GRI = RGBColor(0x33, 0x33, 0x33)
BEYAZ = RGBColor(0xFF, 0xFF, 0xFF)
SARI = RGBColor(0xF3, 0x9C, 0x12)
YESIL = RGBColor(0x27, 0xAE, 0x60)
KIRMIZI = RGBColor(0xE7, 0x4C, 0x3C)
ACIK_MAVI = RGBColor(0xD6, 0xEA, 0xF8)

# --- YARDIMCI FONKSIYONLAR ---

def yeni_sunum():
    """Genis ekran (13.333 x 7.5) bos sunum olustur."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def arkaplan_ekle(slide, renk):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = renk

def baslik_ekle(slide, text, y=Inches(0.3), font_size=40, renk=MAVI_KOYU, bold=True):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = renk
    p.font.bold = bold
    p.font.name = "Calibri"
    return txBox

def alt_baslik_ekle(slide, text, y=Inches(1.3), font_size=18, renk=GRI):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = renk
    p.font.name = "Calibri"

def cizgi_ekle(slide, y=Inches(1.6)):
    line = slide.shapes.add_shape(1, Inches(0.8), y, Inches(11.7), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = MAVI_CANLI
    line.line.fill.background()

def madde_ekle(slide, maddeler, y_bas=Inches(1.8), x=Inches(0.8), font_size=16, renk=GRI, sutun=1):
    if sutun == 1:
        _madde_tek_sutun(slide, maddeler, y_bas, x, font_size, renk)
    else:
        _madde_iki_sutun(slide, maddeler, y_bas, x, font_size, renk)

def _madde_tek_sutun(slide, maddeler, y_bas, x, font_size, renk):
    txBox = slide.shapes.add_textbox(x, y_bas, Inches(11.7), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, madde in enumerate(maddeler):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u25b8 {madde}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = renk
        p.font.name = "Calibri"
        p.space_after = Pt(6)

def _madde_iki_sutun(slide, maddeler, y_bas, x, font_size, renk):
    orta = len(maddeler) // 2
    sol = maddeler[:orta]
    sag = maddeler[orta:]
    
    txBox = slide.shapes.add_textbox(x, y_bas, Inches(5.5), Inches(5.0))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, m in enumerate(sol):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u25b8 {m}"
        p.font.size = Pt(font_size); p.font.color.rgb = renk
        p.font.name = "Calibri"; p.space_after = Pt(5)
    
    txBox2 = slide.shapes.add_textbox(Inches(6.5), y_bas, Inches(6.0), Inches(5.0))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    for i, m in enumerate(sag):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"\u25b8 {m}"
        p.font.size = Pt(font_size); p.font.color.rgb = renk
        p.font.name = "Calibri"; p.space_after = Pt(5)

def tablo_ekle(slide, data, x, y, genislik, yukseklik, baslik_renk=MAVI_KOYU):
    satir = len(data)
    kolon = len(data[0]) if data else 0
    if kolon == 0:
        return
    
    table_shape = slide.shapes.add_table(satir, kolon, x, y, genislik, yukseklik)
    table = table_shape.table
    
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = BEYAZ
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = GRI
                    paragraph.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = baslik_renk
            elif i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ACIK_MAVI

def vurgu_kutusu(slide, metin, y=Inches(4.5), renk=ACIK_MAVI, metin_renk=MAVI_KOYU):
    vurgu = slide.shapes.add_shape(1, Inches(0.8), y, Inches(11.7), Inches(1.5))
    vurgu.fill.solid(); vurgu.fill.fore_color.rgb = renk
    vurgu.line.fill.background()
    tf = vurgu.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = metin
    p.font.size = Pt(16); p.font.color.rgb = metin_renk
    p.font.name = "Calibri"

def kapak_slayt(prs, baslik, alt_baslik, tarih=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    arkaplan_ekle(slide, MAVI_KOYU)
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.5))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = baslik; p.font.size = Pt(60); p.font.color.rgb = SARI
    p.font.bold = True; p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.0))
    tf2 = txBox2.text_frame; p2 = tf2.paragraphs[0]
    p2.text = alt_baslik; p2.font.size = Pt(44); p2.font.color.rgb = BEYAZ
    p2.font.bold = True; p2.font.name = "Calibri"; p2.alignment = PP_ALIGN.CENTER
    
    if tarih:
        txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5))
        tf3 = txBox3.text_frame; p3 = tf3.paragraphs[0]
        p3.text = tarih; p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
        p3.font.name = "Calibri"; p3.alignment = PP_ALIGN.CENTER

def sonuc_slayt(prs, baslik, metin, alt_metin=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    arkaplan_ekle(slide, MAVI_KOYU)
    
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.0))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = baslik; p.font.size = Pt(48); p.font.color.rgb = SARI
    p.font.bold = True; p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.3), Inches(3.0))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = metin
    p2.font.size = Pt(22); p2.font.color.rgb = BEYAZ
    p2.font.name = "Calibri"; p2.alignment = PP_ALIGN.CENTER
    
    if alt_metin:
        p3 = tf2.add_paragraph(); p3.text = alt_metin
        p3.font.size = Pt(20); p3.font.color.rgb = MAVI_ACIK
        p3.font.name = "Calibri"; p3.alignment = PP_ALIGN.CENTER
